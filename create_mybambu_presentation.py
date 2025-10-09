from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Define brand colors
NAVY = RGBColor(20, 29, 56)  # Dark navy background
GREEN = RGBColor(0, 255, 122)  # MyBambu green
WHITE = RGBColor(255, 255, 255)

def add_background(slide, color=NAVY):
    """Add solid color background to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(title_text, subtitle_text=""):
    """Create a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide)

    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2))
    tf = title.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(80)
    p.font.bold = True
    # Split coloring
    if "MyBambu" in title_text:
        parts = title_text.split("MyBambu")
        tf.clear()
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run1 = p.add_run()
        run1.text = "MyBambu"
        run1.font.size = Pt(80)
        run1.font.bold = True
        run1.font.color.rgb = GREEN
        if len(parts) > 1:
            run2 = p.add_run()
            run2.text = parts[1]
            run2.font.size = Pt(80)
            run2.font.bold = True
            run2.font.color.rgb = WHITE
    else:
        p.font.color.rgb = GREEN

    # Subtitle
    if subtitle_text:
        subtitle = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(14), Inches(1))
        stf = subtitle.text_frame
        stf.text = subtitle_text
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sp.font.size = Pt(32)
        sp.font.color.rgb = GREEN

    return slide

# Slide 1: Title
add_title_slide("MyBambu", "Neobank for the US and LATAM")

# Slide 2: Overview
add_title_slide("Overview")

# Slide 3: MyBambu at a Glance
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3)

# Title with mixed colors
title3 = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
tf3 = title3.text_frame
p3 = tf3.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
r1 = p3.add_run()
r1.text = "MyBambu"
r1.font.size = Pt(60)
r1.font.bold = True
r1.font.color.rgb = GREEN
r2 = p3.add_run()
r2.text = " at a Glance"
r2.font.size = Pt(60)
r2.font.bold = True
r2.font.color.rgb = WHITE

# Description
desc3 = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(1.2))
desc3.text_frame.text = "MyBambu's digital modern banking platform best serves all the needs of the unbanked and\nunderbanked Latinos, allowing them transact conveniently and securely in the U.S regardless of their\nimmigration status."
desc3.text_frame.paragraphs[0].font.size = Pt(20)
desc3.text_frame.paragraphs[0].font.color.rgb = GREEN
desc3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Stats boxes
stats = [
    ("15+", "Products\n& Services"),
    ("929,000+", "Customers\nAcquired"),
    ("170+", "Employees"),
    ("40.56m+", "Total\nTransactions")
]

x_start = 1
y_pos = 3.5
box_width = 3
box_height = 2.5
spacing = 0.3

for i, (number, label) in enumerate(stats):
    x_pos = x_start + (i * (box_width + spacing))

    # Box background
    box = slide3.shapes.add_shape(1, Inches(x_pos), Inches(y_pos), Inches(box_width), Inches(box_height))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = GREEN
    box.line.width = Pt(3)

    # Number
    num_box = slide3.shapes.add_textbox(Inches(x_pos), Inches(y_pos + 0.5), Inches(box_width), Inches(1))
    num_box.text_frame.text = number
    num_box.text_frame.paragraphs[0].font.size = Pt(44)
    num_box.text_frame.paragraphs[0].font.bold = True
    num_box.text_frame.paragraphs[0].font.color.rgb = GREEN
    num_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Label
    label_box = slide3.shapes.add_textbox(Inches(x_pos), Inches(y_pos + 1.5), Inches(box_width), Inches(0.8))
    label_box.text_frame.text = label
    label_box.text_frame.paragraphs[0].font.size = Pt(20)
    label_box.text_frame.paragraphs[0].font.bold = True
    label_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 4: Products and Services
add_title_slide("Products and Services")

# Slide 5: Marketplace Differentiation
add_title_slide("Marketplace Differentiation")

# Slide 6: Product Roadmap 2025
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide6)

# Title
title6 = slide6.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
tf6 = title6.text_frame
p6 = tf6.add_paragraph()
p6.alignment = PP_ALIGN.CENTER
r61 = p6.add_run()
r61.text = "Product"
r61.font.size = Pt(60)
r61.font.bold = True
r61.font.color.rgb = GREEN
r62 = p6.add_run()
r62.text = " Roadmap 2025"
r62.font.size = Pt(60)
r62.font.bold = True
r62.font.color.rgb = WHITE

# Roadmap items
roadmap_items = [
    "Loans",
    "Buy Now\nPay Later",
    "Business\nCredit\nCards",
    "IMT\nNetwork\nexpansion",
    "LATAM\nCards",
    "Early ACH\nDeposit\nAccess",
    "P2P & P2B\nReal Time\nPayments",
    "Visa Direct\nDomestic &\nXBorder",
    "Crypto\nWallet\nDomestic &\nXBorder"
]

circle_y = 2.5
circle_size = 1.2
start_x = 0.5
spacing_x = 1.6

for i, item in enumerate(roadmap_items):
    x_pos = start_x + (i * spacing_x)

    # Circle
    circle = slide6.shapes.add_shape(1, Inches(x_pos), Inches(circle_y), Inches(circle_size), Inches(circle_size))
    circle.fill.solid()
    circle.fill.fore_color.rgb = NAVY
    circle.line.color.rgb = GREEN
    circle.line.width = Pt(4)

    # Text
    text_box = slide6.shapes.add_textbox(Inches(x_pos), Inches(circle_y + 0.2), Inches(circle_size), Inches(circle_size - 0.4))
    text_box.text_frame.text = item
    text_box.text_frame.paragraphs[0].font.size = Pt(11)
    text_box.text_frame.paragraphs[0].font.bold = True
    text_box.text_frame.paragraphs[0].font.color.rgb = WHITE
    text_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_box.text_frame.word_wrap = True

# Sub-items for Loans
loan_items = ["Installment\nLoans", "Business\nLoans", "Auto Loans"]
loan_y = 4.5
for i, loan in enumerate(loan_items):
    y_pos = loan_y + (i * 0.8)
    box = slide6.shapes.add_shape(1, Inches(0.3), Inches(y_pos), Inches(1.4), Inches(0.6))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = GREEN

    text = slide6.shapes.add_textbox(Inches(0.3), Inches(y_pos), Inches(1.4), Inches(0.6))
    text.text_frame.text = loan
    text.text_frame.paragraphs[0].font.size = Pt(12)
    text.text_frame.paragraphs[0].font.bold = True
    text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Sub-items for LATAM
latam_items = ["México", "Colombia"]
latam_y = 4.5
for i, country in enumerate(latam_items):
    y_pos = latam_y + (i * 0.8)
    box = slide6.shapes.add_shape(1, Inches(6.7), Inches(y_pos), Inches(1.4), Inches(0.6))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = GREEN

    text = slide6.shapes.add_textbox(Inches(6.7), Inches(y_pos), Inches(1.4), Inches(0.6))
    text.text_frame.text = country
    text.text_frame.paragraphs[0].font.size = Pt(14)
    text.text_frame.paragraphs[0].font.bold = True
    text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Slide 7: NEW - LATAM EXPANSION ROADMAP
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide7)

# Title
title7 = slide7.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
tf7 = title7.text_frame
p7 = tf7.add_paragraph()
p7.alignment = PP_ALIGN.CENTER
r71 = p7.add_run()
r71.text = "LATAM Expansion"
r71.font.size = Pt(60)
r71.font.bold = True
r71.font.color.rgb = GREEN
r72 = p7.add_run()
r72.text = " Roadmap"
r72.font.size = Pt(60)
r72.font.bold = True
r72.font.color.rgb = WHITE

# Subtitle
subtitle7 = slide7.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(0.6))
subtitle7.text_frame.text = "Growing MyBambu's reach across Latin America"
subtitle7.text_frame.paragraphs[0].font.size = Pt(24)
subtitle7.text_frame.paragraphs[0].font.color.rgb = GREEN
subtitle7.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Country expansion boxes
countries = [
    ("Mexico", "Q2 2025"),
    ("Colombia", "Q3 2025"),
    ("Honduras", "Q4 2025"),
    ("Guatemala", "Q1 2026"),
    ("Nicaragua", "Q2 2026"),
    ("El Salvador", "Q3 2026")
]

# Create a visual roadmap with timeline
y_start = 2.8
box_height = 0.9
box_width = 4
spacing_y = 0.15

for i, (country, timeline) in enumerate(countries):
    y_pos = y_start + (i * (box_height + spacing_y))

    # Determine color based on position
    if i < 2:  # Mexico and Colombia - already planned
        box_color = GREEN
        text_color = NAVY
    else:  # New markets
        box_color = WHITE
        text_color = NAVY

    # Country box
    country_box = slide7.shapes.add_shape(1, Inches(2), Inches(y_pos), Inches(box_width), Inches(box_height))
    country_box.fill.solid()
    country_box.fill.fore_color.rgb = box_color
    country_box.line.color.rgb = GREEN
    country_box.line.width = Pt(3)

    # Country name
    country_text = slide7.shapes.add_textbox(Inches(2.2), Inches(y_pos + 0.1), Inches(2.5), Inches(box_height - 0.2))
    country_text.text_frame.text = country
    country_text.text_frame.paragraphs[0].font.size = Pt(28)
    country_text.text_frame.paragraphs[0].font.bold = True
    country_text.text_frame.paragraphs[0].font.color.rgb = text_color
    country_text.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Timeline box
    timeline_box = slide7.shapes.add_shape(1, Inches(6.5), Inches(y_pos), Inches(box_width), Inches(box_height))
    timeline_box.fill.solid()
    timeline_box.fill.fore_color.rgb = NAVY
    timeline_box.line.color.rgb = GREEN
    timeline_box.line.width = Pt(3)

    # Timeline text
    timeline_text = slide7.shapes.add_textbox(Inches(6.7), Inches(y_pos + 0.1), Inches(3.5), Inches(box_height - 0.2))
    timeline_text.text_frame.text = timeline
    timeline_text.text_frame.paragraphs[0].font.size = Pt(24)
    timeline_text.text_frame.paragraphs[0].font.bold = True
    timeline_text.text_frame.paragraphs[0].font.color.rgb = GREEN
    timeline_text.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Key offering text at bottom
offering = slide7.shapes.add_textbox(Inches(1.5), Inches(7.5), Inches(13), Inches(1))
offering.text_frame.text = "Offering: MyBambu digital banking services - mobile money transfers, cash loading, international remittances, and financial inclusion for underserved communities"
offering.text_frame.paragraphs[0].font.size = Pt(16)
offering.text_frame.paragraphs[0].font.color.rgb = WHITE
offering.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
offering.text_frame.word_wrap = True

# Slide 8: Funding
add_title_slide("Funding")

# Slide 9: Capital Invested
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide9)

title9 = slide9.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.5))
tf9 = title9.text_frame
tf9.text = "Bambu Systems LLC\nCapital Invested\nAs of 12/31/2024"
for para in tf9.paragraphs:
    para.alignment = PP_ALIGN.CENTER
    para.font.size = Pt(48)
    para.font.bold = True
    para.font.color.rgb = GREEN

# Investment table
investments = [
    ("Series A", "$53,540,266"),
    ("Series B", "$27,163,000"),
    ("Series C", "$44,107,000"),
    ("Series D", "$68,000,000 (Committed)")
]

table_y = 3
row_height = 0.8
for i, (series, amount) in enumerate(investments):
    y_pos = table_y + (i * row_height)

    # Series label
    label_box = slide9.shapes.add_textbox(Inches(3), Inches(y_pos), Inches(3), Inches(row_height))
    label_box.text_frame.text = series
    label_box.text_frame.paragraphs[0].font.size = Pt(24)
    label_box.text_frame.paragraphs[0].font.color.rgb = WHITE
    label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # Amount
    amount_box = slide9.shapes.add_textbox(Inches(7), Inches(y_pos), Inches(4), Inches(row_height))
    amount_box.text_frame.text = amount
    amount_box.text_frame.paragraphs[0].font.size = Pt(24)
    amount_box.text_frame.paragraphs[0].font.color.rgb = GREEN
    amount_box.text_frame.paragraphs[0].font.bold = True

# Slide 10: Business Performance
add_title_slide("Business Performance")

# Slide 11: All-Time Milestones
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide11)

title11 = slide11.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
tf11 = title11.text_frame
p11 = tf11.add_paragraph()
p11.alignment = PP_ALIGN.CENTER
r111 = p11.add_run()
r111.text = "All-Time Milestones"
r111.font.size = Pt(60)
r111.font.bold = True
r111.font.color.rgb = GREEN
r112 = p11.add_run()
r112.text = " Achieved"
r112.font.size = Pt(60)
r112.font.bold = True
r112.font.color.rgb = WHITE

# Key metrics
metrics = [
    ("$13,300,000+", "Past 12 Months\nRevenue"),
    ("920,000+", "Accounts\nCreated"),
    ("40,560,000+", "Total\nTransactions")
]

metric_y = 3
for i, (number, label) in enumerate(metrics):
    x_pos = 1.5 + (i * 4.5)

    # Number
    num = slide11.shapes.add_textbox(Inches(x_pos), Inches(metric_y), Inches(4), Inches(1))
    num.text_frame.text = number
    num.text_frame.paragraphs[0].font.size = Pt(36)
    num.text_frame.paragraphs[0].font.bold = True
    num.text_frame.paragraphs[0].font.color.rgb = GREEN
    num.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Label
    lbl = slide11.shapes.add_textbox(Inches(x_pos), Inches(metric_y + 1), Inches(4), Inches(1))
    lbl.text_frame.text = label
    lbl.text_frame.paragraphs[0].font.size = Pt(20)
    lbl.text_frame.paragraphs[0].font.color.rgb = WHITE
    lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Transaction breakdown
transactions = [
    ("POS / Card", "$1,716,000,000+"),
    ("Cash Loads", "$583,800,000+"),
    ("ACH", "$580,500,000+"),
    ("Crossborder", "$377,100,000+")
]

trans_y = 5.5
for i, (trans_type, amount) in enumerate(transactions):
    y_pos = trans_y + (i * 0.7)

    # Type
    type_box = slide11.shapes.add_textbox(Inches(9), Inches(y_pos), Inches(3), Inches(0.6))
    type_box.text_frame.text = trans_type
    type_box.text_frame.paragraphs[0].font.size = Pt(18)
    type_box.text_frame.paragraphs[0].font.color.rgb = WHITE

    # Amount
    amt_box = slide11.shapes.add_textbox(Inches(12), Inches(y_pos), Inches(3), Inches(0.6))
    amt_box.text_frame.text = amount
    amt_box.text_frame.paragraphs[0].font.size = Pt(18)
    amt_box.text_frame.paragraphs[0].font.bold = True
    amt_box.text_frame.paragraphs[0].font.color.rgb = GREEN

# Continue with remaining slides (simplified versions)
slide_titles = [
    "Customer Account by Quarter",
    "Historical Transactions",
    "Transactions by Quarter",
    "Transactions by Quarter (Detail)",
    "Customer Revenue and CAC",
    "Customer Revenue and CAC (Detail)",
    "Summarized Figures",
    "Summarized Figures (Detail)",
    "Capitalization Table",
    "MyBambu Board Members",
    "Contact"
]

for title in slide_titles:
    add_title_slide(title)

# Save presentation
output_path = "/Users/vinfa/Desktop/MyBambu_Presentation_with_LATAM_Roadmap.pptx"
prs.save(output_path)
print(f"Presentation created successfully: {output_path}")
